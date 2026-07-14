"""
Dokument-Loader
===============

Extrahiert reinen Text (inkl. Struktur-Metadaten) aus verschiedenen Formaten:

    .pdf          -> PyMuPDF (fitz), Fallback pypdf; Text pro Seite
    .md / .txt    -> direkt (Markdown behält Überschriften-Struktur)
    .docx         -> python-docx
    .pptx         -> python-pptx (Text pro Folie)

Rückgabe: ``LoadedDoc`` mit
    * text   : vollständiger, normalisierter Text
    * blocks : Liste von Textblöcken mit Metadaten (Seite/Folie), damit der
               Chunker Quellenangaben wie "Seite 4" anhängen kann.
"""
from __future__ import annotations

import os
import re
import unicodedata
from dataclasses import dataclass, field
from pathlib import Path

from ragapp.config import settings


def _ocr_log(msg: str) -> None:
    """Freeze-/Crash-SICHERES OCR-Protokoll nach data/logs/ocr.log.

    Jede Zeile wird SOFORT auf die Platte gezwungen (``flush`` + ``os.fsync``) und
    IMMER VOR dem jeweils gefaehrlichen Schritt geschrieben (z. B. vor dem Laden des
    Vision-Modells). Friert das System dabei ein und muss hart neu gestartet werden,
    ist die zuletzt geschriebene Zeile trotzdem erhalten - sie verraet, WOBEI es
    haengen blieb (Modell, Groesse, freier VRAM). Fehler werden verschluckt (das Log
    darf die OCR nie stoeren)."""
    try:
        import time as _t
        from ragapp.config import LOG_DIR
        LOG_DIR.mkdir(parents=True, exist_ok=True)
        line = f"[{_t.strftime('%Y-%m-%d %H:%M:%S')}] {msg}\n"
        with open(LOG_DIR / "ocr.log", "a", encoding="utf-8") as f:
            f.write(line)
            f.flush()
            os.fsync(f.fileno())          # <- erzwingt Schreiben auf die Platte
    except Exception:  # noqa: BLE001
        pass

SUPPORTED_EXTENSIONS = {".pdf", ".md", ".txt", ".markdown", ".docx", ".pptx"}


@dataclass
class Block:
    text: str
    page: int | None = None      # Seiten-/Foliennummer, falls vorhanden
    kind: str = "text"
    ocr: bool = False            # True, wenn der Text per OCR gewonnen wurde
    ocr_engine: str = ""         # "vision" | "easyocr" (nur gesetzt, wenn ocr=True)


@dataclass
class LoadedDoc:
    text: str
    blocks: list[Block]
    filetype: str
    is_markdown: bool = False
    meta: dict = field(default_factory=dict)


# --------------------------------------------------------------------------- #
# Normalisierung
# --------------------------------------------------------------------------- #
def normalize_text(text: str) -> str:
    """Vereinheitlicht Whitespace/Unicode, wichtig für stabile Hashes & Dedup."""
    if not text:
        return ""
    text = unicodedata.normalize("NFC", text)
    text = text.replace("\r\n", "\n").replace("\r", "\n")
    # PDF-Artefakt: nur ECHTE Silbentrennung am Zeilenende zusammenfügen, d. h.
    # Buchstabe + '-' am Zeilenende UND Kleinbuchstabe am nächsten Zeilenanfang
    # (fortgesetztes Wort). Ziffern (Zahlen-/Jahresbereiche wie "2020-\n21") und
    # echte Bindestrich-Komposita ("Nord-\nDeutschland", Fortsetzung großgeschrieben)
    # bleiben so erhalten.
    text = re.sub(r"([A-Za-zÄÖÜäöüß])-\n([a-zäöüß])", r"\1\2", text)
    # überzählige Leerzeichen
    text = re.sub(r"[ \t]+", " ", text)
    # mehr als 2 Leerzeilen -> 2
    text = re.sub(r"\n{3,}", "\n\n", text)
    return text.strip()


# --------------------------------------------------------------------------- #
# Format-spezifische Loader
# --------------------------------------------------------------------------- #
_EASYOCR_READER = None
_EASYOCR_TRIED = False


def _get_easyocr():
    """Gibt einen (gecachten) easyocr-Reader (dt.+engl.) zurueck oder None.

    WICHTIG - laeuft standardmaessig auf der CPU (gpu=False):
    easyocr laedt seine torch-Modelle IN DEN APP-PROZESS. Auf AMD-ROCm ist das
    heikel - eine GPU-Allokation unter VRAM-Druck (Reranker + residentes Chat-
    Modell + evtl. zweite GPU-App) kann den amdgpu-Treiber HART haengen lassen
    (kompletter System-Freeze); ein GPU-Hang ist zudem KEINE Python-Exception, das
    try/except unten faengt ihn also NICHT ab. easyocr ist ohnehin nur der
    Fallback (der Hauptweg ist Vision-OCR ueber Ollama), daher ist CPU hier voll
    ok. Bewusst aktivieren (nur mit reichlich freiem VRAM): RAG_EASYOCR_GPU=1."""
    global _EASYOCR_READER, _EASYOCR_TRIED
    if _EASYOCR_TRIED:
        return _EASYOCR_READER
    _EASYOCR_TRIED = True
    _use_gpu = os.environ.get("RAG_EASYOCR_GPU") == "1"
    try:
        import easyocr
        # CPU-Thread-Deckel: easyocr/torch wuerde sonst ALLE Kerne belegen und den
        # Rechner zusaetzlich ausbremsen. Nur setzen, wenn nicht schon konfiguriert.
        try:
            import torch
            if not _use_gpu and os.environ.get("OMP_NUM_THREADS") is None:
                torch.set_num_threads(max(1, (os.cpu_count() or 4) // 2))
        except Exception:  # noqa: BLE001
            pass
        try:
            _EASYOCR_READER = easyocr.Reader(["de", "en"], gpu=_use_gpu)
        except Exception:  # noqa: BLE001 - GPU evtl. nicht nutzbar -> CPU
            _EASYOCR_READER = easyocr.Reader(["de", "en"], gpu=False)
    except Exception:  # noqa: BLE001 - easyocr nicht installiert
        _EASYOCR_READER = None
    return _EASYOCR_READER


def _easyocr_page(page) -> str:
    """Klassisches OCR einer text-losen Seite: easyocr (GPU, dt.+engl.),
    Fallback pytesseract. '' wenn kein OCR verfuegbar ist."""
    try:
        png = page.get_pixmap(dpi=200).tobytes("png")
    except Exception:  # noqa: BLE001
        return ""
    reader = _get_easyocr()
    if reader is not None:
        try:
            import io
            import numpy as np
            from PIL import Image
            arr = np.array(Image.open(io.BytesIO(png)).convert("RGB"))
            lines = reader.readtext(arr, detail=0, paragraph=True)
            return "\n".join(str(x) for x in lines)
        except Exception:  # noqa: BLE001
            pass
    try:  # Fallback: pytesseract (falls Tesseract installiert ist)
        import io
        import pytesseract
        from PIL import Image
        img = Image.open(io.BytesIO(png))
        try:
            return pytesseract.image_to_string(img, lang="deu") or ""
        except Exception:  # noqa: BLE001
            return pytesseract.image_to_string(img) or ""
    except Exception:  # noqa: BLE001
        return ""


# --------------------------------------------------------------------------- #
# Vision-OCR (Ollama-Multimodal): liest auch Handschrift als zusammenhaengende
# Woerter, wo easyocr nur Zeichenbrei liefert. Nicht garantiert wortgetreu ->
# Decoding wird gegen Wiederholungs-Loops abgesichert, degenerierte Ausgaben
# werden verworfen (Fallback easyocr).
# --------------------------------------------------------------------------- #
_VISION_CAP_CACHE: dict[str, bool] = {}     # model -> vision-faehig?
_VISION_MODEL_RESOLVED = None               # gecachtes Auto-Detektionsergebnis ("" = keins)

# Bevorzugte kleine, laptop-taugliche Vision-Modelle (Reihenfolge = Praeferenz).
_VISION_PREFERRED = ("gemma3:4b", "gemma4:e4b", "gemma3", "gemma4",
                     "minicpm-v", "llava", "qwen2.5vl", "moondream")

_VISION_OCR_PROMPT = (
    "Transkribiere die deutsche Handschrift auf diesem Bild WORTGETREU, "
    "Zeile fuer Zeile, in originaler Reihenfolge. Uebernimm Spiegelstriche (-) "
    "und Pfeile (->). Fasse nichts zusammen, kuerze nichts und erfinde keine "
    "Aufzaehlungen oder Ueberschriften. Ein einzelnes unlesbares Wort ersetzt "
    "du durch [unleserlich]; bei erkennbaren Woertern rate die wahrscheinlichste "
    "Lesart. Wiederhole keine Zeile. Gib NUR den transkribierten Text aus - "
    "ohne Vorspann, ohne Erklaerung, ohne Anfuehrungszeichen."
)

_VISION_PREAMBLE_RE = re.compile(
    r"^\s*(hier ist|hier folgt|die transkription|transkription|"
    r"transkribierter text|text)\b[^\n:]*:\s*",
    re.IGNORECASE,
)


def _model_has_vision(model: str) -> bool:
    """True, wenn das Ollama-Modell die 'vision'-Capability meldet (gecacht)."""
    if not model:
        return False
    if model in _VISION_CAP_CACHE:
        return _VISION_CAP_CACHE[model]
    ok = False
    try:
        import ollama
        info = ollama.Client(host=settings.OLLAMA_BASE_URL, timeout=15).show(model)
        caps = getattr(info, "capabilities", None)
        if caps is None and hasattr(info, "get"):
            caps = info.get("capabilities")
        ok = "vision" in [str(c).lower() for c in (caps or [])]
    except Exception:  # noqa: BLE001
        ok = False
    _VISION_CAP_CACHE[model] = ok
    return ok


def _resolve_vision_model() -> str:
    """Ermittelt das zu nutzende Vision-Modell. Explizit gesetztes
    OCR_VISION_MODEL gewinnt; sonst Auto-Detektion des BESTEN installierten,
    vision-faehigen Modells, das noch in den VRAM passt (statt einfach des
    kleinsten). Massstab ist dasselbe Budget wie bei der Antwort-Modell-Wahl
    (hardware._fit_budget / recommend_ocr_vision_model)."""
    global _VISION_MODEL_RESOLVED
    explicit = (settings.OCR_VISION_MODEL or "").strip()
    if explicit:
        return explicit if _model_has_vision(explicit) else ""
    if _VISION_MODEL_RESOLVED is not None:
        return _VISION_MODEL_RESOLVED
    resolved = ""
    try:
        import ollama
        listing = ollama.Client(host=settings.OLLAMA_BASE_URL, timeout=15).list()
        models = getattr(listing, "models", None)
        if models is None and hasattr(listing, "get"):
            models = listing.get("models", [])
        avail: list[tuple[str, int]] = []
        for m in (models or []):
            name = getattr(m, "model", None) or getattr(m, "name", None)
            if name is None and hasattr(m, "get"):
                name = m.get("model") or m.get("name")
            size = getattr(m, "size", None)
            if size is None and hasattr(m, "get"):
                size = m.get("size")
            if name:
                avail.append((str(name), int(size or 0)))

        # VRAM-Budget + Qualitaets-Reihenfolge aus der Hardware ableiten (einmalig,
        # da _VISION_MODEL_RESOLVED gecacht wird). Robust: faellt auf die alte
        # Heuristik (Praeferenz -> kleinste Datei) zurueck, wenn etwas fehlt.
        fit_bytes, rec_tag, order = 0, "", []
        try:
            from ragapp.hardware import (detect_hardware, _fit_budget,
                                         recommend_ocr_vision_model,
                                         _VISION_OCR_ORDER)
            hw = detect_hardware()
            fit_bytes = int(_fit_budget(hw)[0] * (1024 ** 3))
            rec_tag = (recommend_ocr_vision_model(hw) or "").lower()
            order = list(_VISION_OCR_ORDER)
        except Exception:  # noqa: BLE001
            pass
        # Bewusst KEIN harter Groessen-Cap: das groesste Modell, das in den (Gesamt-)
        # VRAM passt, ist gewollt (z. B. 18 GB auf 24 GB) - es wird nur SICHER
        # nutzbar, weil vor dem OCR alle anderen eigenen Modelle entladen werden und
        # ein freier-VRAM-Gate (_vision_ocr_prepare) prueft, dass es wirklich passt.

        def _qual_rank(name: str) -> int:
            """Kleiner = besser: empfohlenes Modell zuerst, dann Vision-Katalog
            (beste Lesetreue), dann bekannte kleine Vision-Familien."""
            low = name.lower()
            if rec_tag and (low == rec_tag or low.startswith(rec_tag)):
                return -1
            for i, t in enumerate(order):
                if low == t or low.startswith(t):
                    return i
            for j, p in enumerate(_VISION_PREFERRED):
                if low.startswith(p) or p in low:
                    return len(order) + j
            return len(order) + len(_VISION_PREFERRED)

        # Passt komplett in den VRAM? Datei-/Footprint-Groesse als Massstab.
        #   passt   -> beste Qualitaet zuerst, groesseres als Tiebreak
        #   passt nicht -> kleinste zuerst (geringstes Offload)
        def _sort_key(t: tuple[str, int]):
            name, size = t
            if fit_bytes <= 0 or size <= fit_bytes:
                return (0, _qual_rank(name), -size)
            return (1, size, 0)

        avail.sort(key=_sort_key)
        for name, _ in avail:
            if _model_has_vision(name):
                resolved = name
                break
    except Exception:  # noqa: BLE001
        resolved = ""
    _VISION_MODEL_RESOLVED = resolved
    return resolved


def has_vision_ocr_model(pull_if_missing: bool = False,
                         pull_model: str = "") -> str:
    """Public (fuer den Installer): gibt ein installiertes, vision-faehiges Modell
    fuer die Handschrift-/Scan-OCR zurueck (Config ``OCR_VISION_MODEL`` oder
    Auto-Detektion). Ist keins da und ``pull_if_missing=True``, wird das per
    ``hardware.recommend_ocr_vision_model`` empfohlene Vision-Modell gezogen
    (Laptop -> gemma3:4b, mehr VRAM -> gemma3:12b/gemma3:27b). ``pull_model``
    ueberschreibt die Empfehlung explizit. Gibt '' zurueck, wenn keins verfuegbar/
    ziehbar ist (dann faellt die OCR auf easyocr zurueck). Blockiert waehrend des
    Pulls."""
    global _VISION_MODEL_RESOLVED
    m = _resolve_vision_model()
    if m or not pull_if_missing:
        return m
    target = (pull_model or "").strip()
    if not target:
        try:
            from ragapp.hardware import detect_hardware, recommend_ocr_vision_model
            target = recommend_ocr_vision_model(detect_hardware())
        except Exception:  # noqa: BLE001
            target = "gemma3:4b"
    target = target or "gemma3:4b"
    try:
        import ollama
        ollama.Client(host=settings.OLLAMA_BASE_URL, timeout=3600).pull(target)
        _VISION_CAP_CACHE.pop(target, None)
        _VISION_MODEL_RESOLVED = None          # Cache invalidieren -> neu detektieren
        return _resolve_vision_model() or target
    except Exception:  # noqa: BLE001
        return ""


def _render_page_png(page, max_side: int, dpi: int) -> bytes | None:
    """Rendert die Seite und verkleinert sie auf 'max_side' (lange Kante) ->
    weniger VRAM, weniger Wiederholungs-Loops beim Vision-Modell."""
    try:
        import io
        from PIL import Image
        pix = page.get_pixmap(dpi=dpi)
        img = Image.open(io.BytesIO(pix.tobytes("png"))).convert("RGB")
        w, h = img.size
        scale = max_side / max(w, h)
        if scale < 1:
            img = img.resize((max(1, int(w * scale)), max(1, int(h * scale))),
                             Image.LANCZOS)
        buf = io.BytesIO()
        img.save(buf, format="PNG")
        return buf.getvalue()
    except Exception:  # noqa: BLE001
        return None


def _clean_vision_text(raw: str) -> str:
    """Entfernt Code-Fences, Vorspann ('Hier ist die Transkription:') und
    umschliessende Anfuehrungszeichen."""
    t = (raw or "").strip()
    if not t:
        return ""
    if t.startswith("```"):
        t = re.sub(r"^```[a-zA-Z]*\n?", "", t)
        t = re.sub(r"\n?```$", "", t).strip()
    t = _VISION_PREAMBLE_RE.sub("", t, count=1)
    if len(t) >= 2 and t[0] in "\"'" and t[-1] == t[0]:
        t = t[1:-1].strip()
    return t


def _degeneration_kind(text: str) -> "str | None":
    """Art der Vision-Entartung, sonst None. Faengt drei Muster ab, die kleine
    Modelle bei schwerer Handschrift produzieren:
      'repeat'     (a) exakte Zeilen-Wiederholung (eine Zeile dominiert)
                       -> BEHEBBAR: die eindeutigen Zeilen sind meist echt, nur
                          der Loop muss raus (kollabieren + behalten).
      'unreadable' (b) '[unleserlich]'-Flut -> Seite faktisch NICHT gelesen
                       -> UNBRAUCHBAR: nichts Sinnvolles transkribiert.
      'template'   (c) Template-Loop (gleiches Zeilen-Skelett, nur Zahlen wechseln,
                       z. B. 'Aufgabe 1: Lösung / Aufgabe 2: Lösung / …')
                       -> UNBRAUCHBAR: erfundenes Muster, steht so nicht auf der
                          Seite. Kollabieren hilft NICHT (Zeilen unterscheiden sich).
    Nur 'repeat' ist behebbar; (b)/(c) sind Kauderwelsch und duerfen NICHT als
    Vision-Text in den Index (sonst easyocr-Fallback bzw. Seite als unlesbar)."""
    from collections import Counter
    lines = [ln.strip().lower() for ln in text.split("\n") if ln.strip()]
    if len(lines) < 5:
        return None
    n = len(lines)
    # (b) [unleserlich]-Flut ZUERST prüfen -> Seite faktisch nicht gelesen. (Vor der
    # Exakt-Wiederholung, sonst würde eine Flut IDENTISCHER blanker "[unleserlich]"-
    # Zeilen als behebbarer 'repeat' fehlklassifiziert und behalten statt verworfen.)
    if sum(1 for ln in lines if "[unleserlich]" in ln) / n > 0.45:
        return "unreadable"
    # (a) exakte Wiederholung eines ECHTEN Zeileninhalts -> behebbar (kollabieren).
    if Counter(lines).most_common(1)[0][1] / n > 0.6:
        return "repeat"
    # (c) Zeilen-Skelett (Ziffern/Sonderzeichen entfernt) dominiert -> Template-Loop
    def _skel(ln: str) -> str:
        return re.sub(r"\d+", "#", re.sub(r"[^0-9a-zäöüß#]+", " ", ln)).strip()
    skels = Counter(s for s in (_skel(ln) for ln in lines) if s)
    if skels and skels.most_common(1)[0][1] / n > 0.5:
        return "template"
    return None


def _collapse_repeats(text: str, max_repeat: int = 2) -> str:
    """Reduziert direkt aufeinanderfolgende identische Zeilen auf hoechstens
    'max_repeat' (fangt kleinere Loops ab, die den Degenerations-Test knapp
    verfehlen)."""
    out: list[str] = []
    run, last = 0, None
    for line in text.split("\n"):
        key = line.strip().lower()
        if key and key == last:
            run += 1
            if run > max_repeat:
                continue
        else:
            run, last = 1, key
        out.append(line)
    return "\n".join(out).strip()


def _vision_ocr_page(page, model: str) -> str:
    """Rendert die Seite -> verkleinertes PNG -> Ollama-Vision -> bereinigter
    Text. Gibt '' bei Fehler oder degenerierter Ausgabe zurueck (Aufrufer faellt
    dann auf easyocr zurueck)."""
    png = _render_page_png(page, settings.OCR_VISION_MAX_SIDE, settings.OCR_RENDER_DPI)
    if not png:
        return ""
    try:
        import base64
        import ollama
        client = ollama.Client(host=settings.OLLAMA_BASE_URL,
                               timeout=settings.OCR_VISION_TIMEOUT)
        b64 = base64.b64encode(png).decode("ascii")
        resp = client.chat(
            model=model,
            messages=[{"role": "user", "content": _VISION_OCR_PROMPT,
                       "images": [b64]}],
            options={
                "temperature": 0.2, "top_p": 0.9, "top_k": 40,
                "repeat_penalty": 1.3, "repeat_last_n": 64,   # gegen Endlos-Loops
                "num_ctx": 4096,
                "num_predict": settings.OCR_VISION_NUM_PREDICT,
            },
        )
        raw = (resp.get("message", {}) or {}).get("content", "") or ""
    except Exception:  # noqa: BLE001
        return ""
    cleaned = _clean_vision_text(raw)
    if not cleaned:
        return ""                             # kompletter Vision-Fehlschlag -> easyocr
    # [I6] Bei Scan/Handschrift ist Vision meist besser als easyocr - deshalb
    # verwerfen wir einen bloßen WIEDERHOL-Loop nicht komplett, sondern kollabieren
    # ihn und behalten die eindeutigen (echten) Zeilen. ABER: eine '[unleserlich]'-
    # Flut oder ein Template-Loop mit wechselnden Zahlen ist KEIN gelesener Text,
    # sondern erfundenes Muster - das darf NICHT als Vision-Text in den Index
    # (Kollabieren hilft dort nicht, da sich die Zeilen unterscheiden). Solche
    # Seiten geben wir leer zurück -> easyocr-Fallback bzw. Seite bleibt (fast) leer
    # und wird von der Pipeline als 'unvollständig/unlesbar' zur Prüfung markiert.
    kind = _degeneration_kind(cleaned)
    if kind == "repeat":
        return _collapse_repeats(cleaned, max_repeat=1)
    if kind in ("unreadable", "template"):
        return ""
    return _collapse_repeats(cleaned)


# Einmal-pro-Dokument-Entscheidung: darf Vision-OCR (Ollama/GPU) genutzt werden,
# ODER ist zu wenig freier VRAM (dann CPU-easyocr statt GPU-Hang/System-Freeze)?
_VISION_OCR_PREPARED: "bool | None" = None
_OCR_CURRENT_FILE: str = ""              # nur fuers Log (welche Datei wird gerade ge-OCR-t)


def _reset_vision_ocr_prepare() -> None:
    """Vor jedem Dokument aufrufen: die Vision-OCR-Entscheidung neu treffen (der
    freie VRAM kann sich zwischen Laeufen geaendert haben)."""
    global _VISION_OCR_PREPARED
    _VISION_OCR_PREPARED = None


def _vision_ocr_prepare(model: str) -> bool:
    """Bereitet Vision-OCR SICHER vor und entscheidet, ob es genutzt werden darf.

    Genau EINMAL pro Dokument (Ergebnis gecacht bis _reset_vision_ocr_prepare):
      1) Andere EIGENE Ollama-Modelle entladen (z. B. das Chat-Modell) -> nur das
         OCR-Modell soll im VRAM liegen. Fremde Modelle (zweite GPU-App) bleiben (Ro5).
      2) Pruefen, ob das OCR-Modell + Puffer WIRKLICH in den (dann) freien VRAM passt.
         Ollama kennt den vom Reranker / einer zweiten App belegten VRAM NICHT und
         wuerde sonst ueberbuchen -> ROCm/amdgpu-Hang -> kompletter System-Freeze.

    True  -> Vision-OCR erlaubt (Modell passt, nur es liegt im VRAM).
    False -> zu wenig freier VRAM -> Vision-OCR ueberspringen (CPU-easyocr-Fallback).
    Best-effort: laesst sich der VRAM nicht messen, wird NICHT blockiert (True)."""
    global _VISION_OCR_PREPARED
    if _VISION_OCR_PREPARED is not None:
        return _VISION_OCR_PREPARED
    ok = True
    _ocr_log(f"=== OCR-Lauf fuer {_OCR_CURRENT_FILE or '?'} (OCR_ENGINE={settings.OCR_ENGINE}) ===")
    try:
        from ragapp import hardware
        gpu = hardware.detect_gpu() or {}
        _ocr_log(f"prepare: GPU={gpu.get('name')!r} vram_total={gpu.get('vram_gb')}GB "
                 f"is_igpu={gpu.get('is_igpu')} | OCR-Modell={model!r}")
        # Kein dedizierter GPU-Speicher (iGPU/CPU-only): Ollama nutzt RAM, kein
        # GPU-Hang-Risiko -> Vision-OCR ohne Sonderbehandlung erlauben.
        if gpu.get("is_igpu") or not gpu.get("vram_gb"):
            _ocr_log("prepare: kein dedizierter VRAM -> Vision-OCR ohne Gate erlaubt")
            _VISION_OCR_PREPARED = True
            return True
        _ocr_log(f"prepare: freier VRAM VOR Entladen={hardware.vram_free_gb()}GB")
        # 1) Andere eigene Modelle entladen -> nur EIN Modell fuer die OCR-Aufgabe.
        try:
            from ragapp.scripts.stop_ollama_standby import unload_resident_models
            _n = unload_resident_models(settings.OLLAMA_BASE_URL)
            _ocr_log(f"prepare: {_n} eigene(s) Ollama-Modell(e) entladen")
            if _n:
                import time
                time.sleep(1.0)   # kurz warten, bis der VRAM wirklich frei gemessen wird
        except Exception as _exc:  # noqa: BLE001
            _ocr_log(f"prepare: Entladen fehlgeschlagen: {_exc!r}")
        # 2) Passt das OCR-Modell jetzt in den ECHTEN freien VRAM (inkl. Reranker /
        #    fremder GPU-App, die Ollama nicht kennt)?
        free = hardware.vram_free_gb()
        need = None
        try:
            from ragapp.llm import _model_size_gb
            need = _model_size_gb(model)
        except Exception:  # noqa: BLE001
            need = None
        headroom = float(getattr(settings, "OCR_VISION_VRAM_HEADROOM_GB", 2.0) or 2.0)
        if free is not None and need and free < need + headroom:
            ok = False
        _ocr_log(f"prepare: freier VRAM NACH Entladen={free}GB, Modell braucht~{need}GB "
                 f"(+{headroom}GB Puffer) -> "
                 f"{'ERLAUBT: Vision-OCR' if ok else 'ZU WENIG VRAM -> CPU-Fallback (kein Vision)'}")
    except Exception as _exc:  # noqa: BLE001
        _ocr_log(f"prepare: Fehler ({_exc!r}) -> best-effort erlaubt")
        ok = True
    _VISION_OCR_PREPARED = ok
    return ok


def _ocr_page(page) -> tuple[str, str]:
    """OCR einer text-losen PDF-Seite. Waehlt die Engine gemaess
    settings.OCR_ENGINE ('vision' | 'easyocr' | 'auto'). 'auto' = Vision, falls
    ein vision-faehiges Ollama-Modell installiert ist, sonst easyocr; bei
    Vision-Fehler/Loop Fallback auf easyocr.
    Rueckgabe: (text, engine) - engine in {'vision', 'easyocr', ''}."""
    engine = (settings.OCR_ENGINE or "auto").strip().lower()
    _pno = getattr(page, "number", None)
    _pno = (_pno + 1) if isinstance(_pno, int) else "?"
    if engine in ("vision", "auto"):
        model = _resolve_vision_model()
        # _vision_ocr_prepare: entlaedt andere Modelle + prueft, ob das OCR-Modell
        # SICHER in den freien VRAM passt. Passt es nicht (z. B. zweite GPU-App belegt
        # VRAM), wird Vision-OCR uebersprungen -> KEIN GPU-Ueberlauf/Freeze, stattdessen
        # CPU-easyocr weiter unten.
        if model and _vision_ocr_prepare(model):
            try:
                from ragapp import hardware
                _fv = hardware.vram_free_gb()
            except Exception:  # noqa: BLE001
                _fv = None
            # DIESE Zeile wird per fsync auf die Platte gezwungen, BEVOR das Modell
            # geladen wird -> friert es hier ein, ist sie nach dem Neustart da.
            _ocr_log(f"Seite {_pno}: >>> Vision-OCR START (Ollama) Modell={model!r} "
                     f"freierVRAM={_fv}GB -- laedt jetzt das Modell/liest die Seite")
            txt = _vision_ocr_page(page, model)
            _ocr_log(f"Seite {_pno}: <<< Vision-OCR fertig (len={len(txt.strip())})")
            if txt.strip():
                return txt, "vision"
            # nur bei komplettem Vision-Fehlschlag (leer) -> easyocr, damit
            # ueberhaupt Text entsteht (degeneriertes Vision wird oben behalten)
    _ocr_log(f"Seite {_pno}: >>> easyocr (CPU) START")
    txt = _easyocr_page(page)
    _ocr_log(f"Seite {_pno}: <<< easyocr fertig (len={len(txt.strip())})")
    return (txt, "easyocr" if txt.strip() else "")


def _load_pdf(path: Path, progress=None) -> LoadedDoc:
    # Vision-OCR-Entscheidung pro Dokument neu treffen (freier VRAM kann sich
    # geaendert haben); tatsaechlich entladen/geprueft wird erst beim ERSTEN OCR-Bedarf.
    global _OCR_CURRENT_FILE
    _OCR_CURRENT_FILE = getattr(path, "name", str(path))
    _reset_vision_ocr_prepare()
    blocks: list[Block] = []
    text_parts: list[str] = []
    ocr_pages: list[int] = []
    ocr_engines: set[str] = set()
    ocr_low_pages = 0        # F2: OCR-Seiten mit < OCR_MIN_PAGE_CHARS Zeichen (unvollstaendig gelesen)
    try:
        import fitz  # PyMuPDF

        with fitz.open(str(path)) as doc:
            n_pages = doc.page_count
            for i, page in enumerate(doc, start=1):
                # sort=True: Lesereihenfolge nach Layout (oben->unten, links->rechts)
                # statt Content-Stream-Reihenfolge -> korrekt bei mehrspaltigen
                # Seiten/Folien.
                raw = page.get_text("text", sort=True) or ""
                norm = normalize_text(raw)
                ocr_engine = ""
                if not norm:
                    # Keine Textebene -> vermutlich Scan/Bild -> OCR versuchen.
                    if progress:
                        try:
                            progress(f"OCR Seite {i}/{n_pages} …", i, n_pages)
                        except Exception:  # noqa: BLE001
                            pass
                    ocr_text, ocr_engine = _ocr_page(page)
                    norm = normalize_text(ocr_text)
                    # F2: sehr wenig OCR-Text -> Seite unvollstaendig gelesen (zaehlt
                    # AUCH die leer gebliebene Seite, norm == "").
                    if len(norm) < settings.OCR_MIN_PAGE_CHARS:
                        ocr_low_pages += 1
                if norm:
                    blk = Block(text=norm, page=i, kind="page")
                    if ocr_engine:
                        blk.ocr = True
                        blk.ocr_engine = ocr_engine
                        ocr_pages.append(i)
                        ocr_engines.add(ocr_engine)
                    blocks.append(blk)
                    text_parts.append(norm)
    except Exception as exc:  # Fallback auf pypdf
        # evtl. schon teilbefüllte fitz-Ergebnisse verwerfen -> keine Seiten-Duplikate
        blocks, text_parts = [], []
        ocr_pages, ocr_engines = [], set()
        ocr_low_pages = 0
        try:
            from pypdf import PdfReader

            reader = PdfReader(str(path))
            for i, page in enumerate(reader.pages, start=1):
                norm = normalize_text(page.extract_text() or "")
                if norm:
                    blocks.append(Block(text=norm, page=i, kind="page"))
                    text_parts.append(norm)
        except Exception as exc2:
            raise RuntimeError(f"PDF konnte nicht gelesen werden: {exc} / {exc2}")
    return LoadedDoc(
        text="\n\n".join(text_parts),
        blocks=blocks,
        filetype="pdf",
        is_markdown=False,
        meta={
            "pages": len(blocks),
            "ocr": bool(ocr_pages),
            "ocr_pages": ocr_pages,
            "ocr_engine": ",".join(sorted(ocr_engines)),
            "ocr_low_pages": ocr_low_pages,   # F2: Anzahl unvollstaendig gelesener OCR-Seiten
        },
    )


def _load_markdown(path: Path) -> LoadedDoc:
    raw = path.read_text(encoding="utf-8", errors="replace")
    # optionales YAML-Frontmatter entfernen (Metadaten)
    fm_meta: dict = {}
    try:
        import frontmatter

        post = frontmatter.loads(raw)
        fm_meta = dict(post.metadata)
        raw = post.content
    except Exception:
        pass
    norm = normalize_text(raw)
    return LoadedDoc(
        text=norm,
        blocks=[Block(text=norm, page=None, kind="markdown")],
        filetype="md",
        is_markdown=True,
        meta=fm_meta,
    )


def _load_txt(path: Path) -> LoadedDoc:
    norm = normalize_text(path.read_text(encoding="utf-8", errors="replace"))
    return LoadedDoc(
        text=norm, blocks=[Block(text=norm, kind="text")], filetype="txt"
    )


def _table_to_md(rows: list) -> str:
    """Wandelt eine Tabelle (Liste von Zeilen mit Zellen) in eine Markdown-Tabelle -
    so bleiben Zeilen/Spalten-Beziehungen (Pro/Contra, Klassifikationen, Kennzahlen)
    erhalten, statt zu Zeichenbrei zu kollabieren."""
    clean = [[str(c).replace("\n", " ").replace("|", "/").strip() for c in r]
             for r in rows if any(str(c).strip() for c in r)]
    if not clean:
        return ""
    ncol = max(len(r) for r in clean)
    clean = [r + [""] * (ncol - len(r)) for r in clean]
    out = ["| " + " | ".join(clean[0]) + " |",
           "| " + " | ".join(["---"] * ncol) + " |"]
    for r in clean[1:]:
        out.append("| " + " | ".join(r) + " |")
    return "\n".join(out)


_DOCX_HEADING_RE = re.compile(r"^(?:heading|überschrift|ueberschrift)\s*(\d+)$")


def _docx_heading_level(style_name: str | None) -> int | None:
    """Ermittelt die Markdown-Überschriftenebene (1..6) aus einem DOCX-Absatzstil.
    'Heading 1'..'Heading 9' -> 1..6 (auf 6 gedeckelt), 'Title'/'Titel' -> 1,
    'Subtitle'/'Untertitel' -> 2. Sonst None (normaler Absatz). Reine Funktion,
    damit sie ohne python-docx testbar ist."""
    if not style_name:
        return None
    s = style_name.strip().lower()
    m = _DOCX_HEADING_RE.match(s)
    if m:
        return min(max(int(m.group(1)), 1), 6)
    if s in ("title", "titel"):
        return 1
    if s in ("subtitle", "untertitel"):
        return 2
    return None


def _load_docx(path: Path) -> LoadedDoc:
    from docx import Document

    doc = Document(str(path))
    # Überschriften-Stile (Heading 1/2/3 …) als Markdown-Struktur (# / ## / ###)
    # rekonstruieren -> is_markdown=True lässt das semantische Chunking greifen und
    # macht 'location' brauchbar (statt eines einzigen 'Dokument'-Blocks).
    lines: list[str] = []
    for p in doc.paragraphs:
        text = p.text.strip()
        if not text:
            continue
        try:
            style_name = p.style.name if p.style is not None else None
        except Exception:  # noqa: BLE001 - defekter/fehlender Stilverweis
            style_name = None
        level = _docx_heading_level(style_name)
        if level:
            lines.append("#" * level + " " + text)
        else:
            lines.append(text)
    # Tabellen als Markdown erfassen (Struktur bleibt erhalten)
    for table in doc.tables:
        md = _table_to_md([[c.text for c in row.cells] for row in table.rows])
        if md:
            lines.append(md)
    norm = normalize_text("\n".join(lines))
    return LoadedDoc(
        text=norm,
        blocks=[Block(text=norm, page=None, kind="markdown")],
        filetype="docx",
        is_markdown=True,
    )


def _iter_pptx_shapes(shapes):
    """Iteriert Folien-Shapes rekursiv - auch in Gruppen/SmartArt (die sonst
    komplett verloren gingen)."""
    for shape in shapes:
        yield shape
        if hasattr(shape, "shapes"):          # GroupShape
            try:
                yield from _iter_pptx_shapes(shape.shapes)
            except Exception:  # noqa: BLE001
                pass


def _load_pptx(path: Path) -> LoadedDoc:
    from pptx import Presentation

    prs = Presentation(str(path))
    blocks: list[Block] = []
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, start=1):
        texts: list[str] = []
        for shape in _iter_pptx_shapes(slide.shapes):
            if getattr(shape, "has_text_frame", False) and shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = "".join(run.text for run in para.runs)
                    if line.strip():
                        texts.append(line)
            if getattr(shape, "has_table", False) and shape.has_table:
                md = _table_to_md([[c.text for c in row.cells] for row in shape.table.rows])
                if md:
                    texts.append(md)
        # Sprechernotizen: oft die AUSFORMULIERTE Erklaerung, die die knappe Folie
        # nur andeutet - bisher komplett verworfen.
        try:
            if slide.has_notes_slide:
                notes = (slide.notes_slide.notes_text_frame.text or "").strip()
                if notes:
                    texts.append("Notizen: " + notes)
        except Exception:  # noqa: BLE001
            pass
        norm = normalize_text("\n".join(texts))
        if norm:
            blocks.append(Block(text=norm, page=i, kind="slide"))
            parts.append(norm)
    return LoadedDoc(
        text="\n\n".join(parts), blocks=blocks, filetype="pptx",
        meta={"slides": len(blocks)},
    )


_LOADERS = {
    ".pdf": _load_pdf,
    ".md": _load_markdown,
    ".markdown": _load_markdown,
    ".txt": _load_txt,
    ".docx": _load_docx,
    ".pptx": _load_pptx,
}


def load_document(path: str | Path, progress=None) -> LoadedDoc:
    path = Path(path)
    ext = path.suffix.lower()
    if ext not in _LOADERS:
        raise ValueError(f"Nicht unterstütztes Format: {ext} ({path.name})")
    # Nur der PDF-Loader kann Seiten rendern -> OCR-Fortschritt melden.
    if ext == ".pdf":
        return _LOADERS[ext](path, progress=progress)
    return _LOADERS[ext](path)
